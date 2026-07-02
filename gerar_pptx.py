#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Emu(12192000)  # 16:9
prs.slide_height = Emu(6858000)

BLUE = RGBColor(0x15, 0x65, 0xC0)
DARK_BLUE = RGBColor(0x0D, 0x47, 0x8A)
DARK = RGBColor(0x1F, 0x29, 0x37)
GRAY = RGBColor(0x6B, 0x72, 0x78)
CREAM = RGBColor(0xF7, 0xF5, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)

def set_font(run, size, bold=False, italic=False, color=DARK, font_name='Calibri'):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name

def add_bg(slide, color=CREAM):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_eyebrow(slide, text, left=Inches(0.8), top=Inches(0.4)):
    txBox = slide.shapes.add_textbox(left, top, Inches(10), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_font(run, 11, bold=True, color=BLUE)
    p.alignment = PP_ALIGN.LEFT

def add_title(slide, text, left=Inches(0.8), top=Inches(0.9), size=44):
    txBox = slide.shapes.add_textbox(left, top, Inches(12), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_font(run, size, bold=True, color=DARK)

def add_body(slide, text, left=Inches(0.8), top=Inches(2.2), width=Inches(10), size=16):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_font(run, size, color=DARK)

def add_card_rect(slide, left, top, width, height, fill=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

# ── SLIDE 1: CAPA ──
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

add_eyebrow(slide, "GRUPO 3 · TRIBUNAL DAS COMUNIDADES TRADICIONAIS")

add_title(slide, "Comunidades\nRibeirinhas", top=Inches(1.2), size=52)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(9), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = '"O rio é nossa casa, nossa fonte de vida e nossa identidade."'
set_font(run, 22, italic=True, color=DARK_BLUE)

# Cards
cards = [
    ("01", "Quem somos", "Famílias que vivem às margens do Rio Água Boa, entre pesca artesanal, extrativismo e agricultura de várzea, há gerações."),
    ("02", "Nosso território", "O rio é estrada, alimento e identidade — o ciclo de cheia e vazante organiza toda a nossa vida e nossa economia."),
    ("03", "O que defendemos", "Consulta prévia, estudo de impacto ambiental independente e preservação do rio antes de qualquer decisão.")
]
for i, (num, title, text) in enumerate(cards):
    x = Inches(0.8 + i * 4.2)
    y = Inches(4.0)
    card = add_card_rect(slide, x, y, Inches(3.8), Inches(2.0))
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(18)
    tf.margin_right = Pt(18)
    tf.margin_top = Pt(14)
    tf.margin_bottom = Pt(14)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = num
    set_font(run, 24, bold=True, color=BLUE)
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = title
    set_font(run2, 14, bold=True, color=DARK)
    p3 = tf.add_paragraph()
    run3 = p3.add_run()
    run3.text = text
    set_font(run3, 12, color=GRAY)

# Footer
add_card_rect(slide, Inches(0.8), Inches(6.5), Inches(14), Inches(0.01), fill=RGBColor(0xD9,0xD9,0xD3))
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(13), Inches(0.6))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "SITUAÇÃO-PROBLEMA"
set_font(run, 10, bold=True, color=GRAY)
p2 = tf.add_paragraph()
run2 = p2.add_run()
run2.text = "Um empreendimento pretende se instalar no território do Rio Água Boa, ameaçando o equilíbrio ambiental e o modo de vida ribeirinho."
set_font(run2, 13, color=DARK)

# ── SLIDES 2-6: CONTEÚDO ──
slides_data = [
    ("PERGUNTA", "Quem somos?", "RESPOSTA", [
        "Somos famílias ribeirinhas que habitam as margens do Rio Água Boa há mais de um século. Nossa comunidade é formada por pescadores artesanais, extrativistas, agricultores de várzea e artesãos que dependem diretamente dos recursos do rio para viver.",
        "Nossa identidade é construída na relação de respeito e pertencimento com as águas. O saber tradicional é transmitido de geração em geração — os mais velhos ensinam a ler os sinais da natureza, os períodos de cheia e vazante, e as técnicas de pesca que não esgotam os recursos.",
        "Não somos moradores da beira do rio: somos o próprio rio em forma de gente."
    ], None),
    ("PERGUNTA", "Como vivemos?", "RESPOSTA", [
        "Nosso modo de vida é organizado pelo ritmo do rio. Na cheia, os peixes se reproduzem e as áreas de várzea são fertilizadas; na vazante, plantamos e colhemos. Cada estação tem seu trabalho, sua comida e sua festa.",
    ], [
        ("Pesca artesanal", "peixes como curimatã, tambaqui e tucunaré sustentam nossa alimentação e geram renda nas feiras da cidade."),
        ("Agricultura de várzea", "plantamos mandioca, milho, feijão e frutas nas terras férteis que o rio renova a cada enchente."),
        ("Extrativismo", "coletamos açaí, borracha, castanha e plantas medicinais da floresta alagável."),
        ("Artesanato", "com cipó, palha e barro, criamos peças que carregam nossa cultura e nossa história."),
    ]),
    ("PERGUNTA", "Qual a importância do território para nós?", "RESPOSTA", [
        "O território não é apenas o lugar onde moramos — é a condição de existência da nossa comunidade. Sem o Rio Água Boa em equilíbrio, nossa cultura, nossa economia e nossa identidade simplesmente deixam de existir.",
    ], [
        ("Identidade cultural", "as narrativas, os saberes, as festas e as crenças da nossa comunidade nascem da relação com o rio."),
        ("Sobrevivência material", "é do território que tiramos nossa comida, nossa água, nossa moradia e nossa renda."),
        ("Território como sujeito", "para nós, o rio não é um recurso a ser explorado — é um ente vivo com o qual nos relacionamos."),
        ("Direito constitucional", "a Constituição Federal de 1988 garante às comunidades tradicionais o direito ao território."),
    ]),
    ("PERGUNTA", "Benefícios e prejuízos do projeto?", None, [], None),
    ("PERGUNTA", "O que defendemos?", None, [], [
        ("Consulta prévia, livre e informada", "Convenção 169 da OIT, da qual o Brasil é signatário."),
        ("EIA independente", "Estudo de Impacto Ambiental por equipe escolhida em acordo com a comunidade."),
        ("Preservação do Rio Água Boa", "ciclo hidrológico, qualidade da água e estoques pesqueiros."),
        ("Participação nos benefícios", "compensação justa e participação nos resultados se o projeto for viável."),
        ("Nenhuma remoção forçada", "o território ribeirinho é nossa casa e não será negociado sob pressão."),
        ("Fiscalização permanente", "participação da comunidade no monitoramento ambiental contínuo."),
    ]),
]

for idx, (qlabel, qtext, rlabel, paragraphs, list_items) in enumerate(slides_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    n = idx + 2
    add_eyebrow(slide, "COMUNIDADES RIBEIRINHAS · TRIBUNAL DAS COMUNIDADES TRADICIONAIS")

    # Question card
    QC_TOP = Inches(0.9)
    QC_LEFT = Inches(0.8)
    card = add_card_rect(slide, QC_LEFT, QC_TOP, Inches(13), Inches(1.2))
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(24)
    tf.margin_right = Pt(24)
    tf.margin_top = Pt(16)
    tf.margin_bottom = Pt(16)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = qlabel
    set_font(run, 11, bold=True, color=BLUE)
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = qtext
    set_font(run2, 36, bold=True, color=DARK)

    if n == 2:
        # Single answer card
        AC_TOP = Inches(2.4)
        card = add_card_rect(slide, Inches(0.8), AC_TOP, Inches(13), Inches(2.8))
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(24)
        tf.margin_right = Pt(24)
        tf.margin_top = Pt(18)
        tf.margin_bottom = Pt(14)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "RESPOSTA"
        set_font(run, 11, bold=True, color=GRAY)
        for para in paragraphs:
            p2 = tf.add_paragraph()
            run2 = p2.add_run()
            run2.text = para
            set_font(run2, 16, color=DARK if ":" in para else DARK)

    elif n in (3, 4):
        # Answer card + split list
        AC_TOP = Inches(2.4)
        card = add_card_rect(slide, Inches(0.8), AC_TOP, Inches(13), Inches(1.2))
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(24)
        tf.margin_right = Pt(24)
        tf.margin_top = Pt(14)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "RESPOSTA"
        set_font(run, 11, bold=True, color=GRAY)
        for para in paragraphs:
            p2 = tf.add_paragraph()
            run2 = p2.add_run()
            run2.text = para
            set_font(run2, 16, color=DARK)

        # List cards
        COL_W = Inches(6.2)
        items_left = list_items[:2]
        items_right = list_items[2:]
        for col_idx, items in enumerate([items_left, items_right]):
            x = Inches(0.8 + col_idx * 6.8)
            card = add_card_rect(slide, x, Inches(3.9), COL_W, Inches(2.0))
            tf = card.text_frame
            tf.word_wrap = True
            tf.margin_left = Pt(20)
            tf.margin_right = Pt(16)
            tf.margin_top = Pt(14)
            tf.margin_bottom = Pt(8)
            for i, (title, desc) in enumerate(items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                run = p.add_run()
                run.text = f"{title}: "
                set_font(run, 15, bold=True, color=BLUE)
                run2 = p.add_run()
                run2.text = desc
                set_font(run2, 14, color=DARK)

    elif n == 5:
        # Pro/Con
        COL_W = Inches(6.2)
        pro_items = [
            "Geração de empregos temporários na construção civil e serviços locais.",
            "Aumento da arrecadação municipal e promessa de investimentos em infraestrutura.",
            "Demanda por produtos e serviços locais durante a operação.",
        ]
        con_items = [
            "Poluição do Rio Água Boa por rejeitos industriais, comprometendo a pesca e o abastecimento de água.",
            "Deslocamento forçado de famílias e destruição de áreas de várzea férteis para plantio.",
            "Rompimento dos laços comunitários e da transmissão de saberes tradicionais.",
            "Empregos temporários e de baixa qualificação — o lucro vai para acionistas, o passivo ambiental fica conosco.",
        ]
        for col_idx, (items, color, title) in enumerate([
            (pro_items, GREEN, "Possíveis benefícios (segundo a empresa)"),
            (con_items, RED, "Prejuízos comprovados (nossa perspectiva)")
        ]):
            x = Inches(0.8 + col_idx * 6.8)
            card = add_card_rect(slide, x, Inches(2.4), COL_W, Inches(3.8))
            tf = card.text_frame
            tf.word_wrap = True
            tf.margin_left = Pt(20)
            tf.margin_right = Pt(16)
            tf.margin_top = Pt(18)
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title
            set_font(run, 16, bold=True, color=color)
            for item in items:
                p2 = tf.add_paragraph()
                run2 = p2.add_run()
                run2.text = f"→ {item}"
                set_font(run2, 14, color=DARK)

    elif n == 6:
        # Highlight + 2 col list
        HL_TOP = Inches(2.4)
        card = add_card_rect(slide, Inches(0.8), HL_TOP, Inches(13), Inches(0.9))
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(24)
        tf.margin_right = Pt(24)
        tf.margin_top = Pt(16)
        tf.margin_bottom = Pt(10)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = '"Não somos contra o desenvolvimento, mas contra um desenvolvimento que nos exclui e destrói o que somos."'
        set_font(run, 17, italic=True, color=BLUE)

        COL_W = Inches(6.2)
        items_left = list_items[:3]
        items_right = list_items[3:]
        for col_idx, items in enumerate([items_left, items_right]):
            x = Inches(0.8 + col_idx * 6.8)
            card = add_card_rect(slide, x, Inches(3.6), COL_W, Inches(2.6))
            tf = card.text_frame
            tf.word_wrap = True
            tf.margin_left = Pt(20)
            tf.margin_right = Pt(16)
            tf.margin_top = Pt(14)
            for i, (title, desc) in enumerate(items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                run = p.add_run()
                run.text = f"{title}: "
                set_font(run, 15, bold=True, color=BLUE)
                run2 = p.add_run()
                run2.text = desc
                set_font(run2, 14, color=DARK)

    # Slide number
    txBox = slide.shapes.add_textbox(Inches(13), Inches(7.0), Inches(1.2), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{n} / 6"
    set_font(run, 12, bold=True, color=BLUE)

prs.save('/root/slides-completo.pptx')
print("PPTX saved: /root/slides-completo.pptx")
